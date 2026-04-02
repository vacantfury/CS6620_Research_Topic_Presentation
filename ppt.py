"""
Generate a PowerPoint presentation for CS 6620 Research Topic:
"Serverless Computing Meets LLM Inference: Current Research and Future Directions"

Focused on 3 core papers:
  1. Hassan et al. (2021) - Survey on Serverless Computing [assigned]
  2. Fu et al. (2024) - ServerlessLLM (OSDI '24)
  3. Wen et al. (2026) - SlsDetector (ACM TOSEM)

With brief mentions of 3 supplementary papers where relevant.

Usage: python ppt.py
Output: presentation.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Color Palette ───────────────────────────────────────────────────────────
DARK_NAVY = RGBColor(0x0B, 0x1D, 0x3A)
DARK_BLUE = RGBColor(0x14, 0x2D, 0x5E)
MEDIUM_BLUE = RGBColor(0x1E, 0x40, 0x7A)
TEAL = RGBColor(0x0E, 0x8C, 0x8C)
LIGHT_TEAL = RGBColor(0x14, 0xB8, 0xA6)
CORAL = RGBColor(0xF4, 0x6D, 0x5B)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
MEDIUM_GRAY = RGBColor(0x94, 0xA3, 0xB8)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
MUTED_TEXT = RGBColor(0x64, 0x74, 0x8B)

# ─── Paths ───────────────────────────────────────────────────────────────────
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
FIGURES_DIR = os.path.join(SCRIPT_DIR, "figures")
OUTPUT_FILE = os.path.join(SCRIPT_DIR, "presentation.pptx")

# ─── Slide dimensions (16:9 widescreen) ─────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ═══════════════════════════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ═══════════════════════════════════════════════════════════════════════════════

def set_slide_bg(slide, color):
    """Set the background color of a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None):
    """Add a shape to the slide."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri", line_spacing=1.2):
    """Add a text box with a single paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK_TEXT, font_name="Calibri", bullet_color=None,
                    spacing=1.4):
    """Add a bullet list. Items can be strings or (bold_part, rest) tuples."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * spacing)

        if isinstance(item, tuple):
            run1 = p.add_run()
            run1.text = "▸ " + item[0]
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = bullet_color or TEAL
            run1.font.bold = True
            run1.font.name = font_name

            run2 = p.add_run()
            run2.text = item[1]
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.bold = False
            run2.font.name = font_name
        else:
            run = p.add_run()
            run.text = "▸ " + item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = False
            run.font.name = font_name

    return txBox


def add_accent_bar(slide, left, top, width, height, color=TEAL):
    """Add a thin accent bar."""
    return add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height, fill_color=color)


def add_section_number(slide, number, left, top, size=Inches(0.7)):
    """Add a circled section number."""
    shape = add_shape(slide, MSO_SHAPE.OVAL, left, top, size, size, fill_color=TEAL)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = str(number)
    run.font.size = Pt(24)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = "Calibri"
    return shape


def add_card(slide, left, top, width, height, fill_color=WHITE):
    """Add a rounded rectangle card."""
    shape = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
                      fill_color=fill_color)
    shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape


def add_header_bar(slide, title, section_num=None):
    """Add the standard dark header bar with title."""
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), fill_color=DARK_NAVY)
    if section_num is not None:
        add_section_number(slide, section_num, Inches(0.6), Inches(0.2), Inches(0.55))
        add_textbox(slide, Inches(1.4), Inches(0.25), Inches(11), Inches(0.6),
                    title, font_size=28, color=WHITE, bold=True, font_name="Calibri")
    else:
        add_textbox(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.6),
                    title, font_size=28, color=WHITE, bold=True, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS — 9 slides total
# ═══════════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)

    # Top & bottom accent bars
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), fill_color=TEAL)
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06), fill_color=TEAL)

    # Left accent line
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(1.0), Inches(1.8), Inches(0.06), Inches(3.8), fill_color=TEAL)

    # Title text
    add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.0),
                "Serverless Computing Meets", font_size=42, color=LIGHT_TEAL,
                bold=True)
    add_textbox(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(1.0),
                "LLM Inference", font_size=52, color=WHITE, bold=True)

    # Subtitle
    add_textbox(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.8),
                "Current Research and Future Directions", font_size=24,
                color=MEDIUM_GRAY)

    # Presenter names
    add_textbox(slide, Inches(1.5), Inches(4.6), Inches(10), Inches(0.5),
                "Haoyu Zhang  &  Esme Wang", font_size=20, color=WHITE)

    # Divider
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(1.5), Inches(5.2), Inches(3), Inches(0.03), fill_color=TEAL)

    # Course info
    add_textbox(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
                "CS 6620 — Cloud Computing", font_size=18, color=MEDIUM_GRAY)
    add_textbox(slide, Inches(1.5), Inches(6.0), Inches(10), Inches(0.5),
                "Research Topic Presentation  •  Spring 2026", font_size=16,
                color=MUTED_TEXT)


def slide_02_background(prs):
    """Slide 2: What is Serverless Computing? (Background)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_header_bar(slide, "WHAT IS SERVERLESS COMPUTING?", section_num=1)

    # Architecture diagram on the left
    arch_path = os.path.join(FIGURES_DIR, "serverless_architecture.png")
    if os.path.exists(arch_path):
        slide.shapes.add_picture(arch_path,
                                 Inches(0.4), Inches(1.3), Inches(5.8), Inches(5.8))

    # Key characteristics on the right
    add_textbox(slide, Inches(6.7), Inches(1.4), Inches(6), Inches(0.5),
                "Key Characteristics", font_size=24, color=DARK_NAVY, bold=True)
    add_accent_bar(slide, Inches(6.7), Inches(1.95), Inches(2.0), Inches(0.04))

    items = [
        ("Auto-scaling: ", "Applications scale up/down on demand automatically"),
        ("Pay-per-use: ", "Charged only when your function is actually running"),
        ("No server management: ", "Cloud provider handles all infrastructure"),
        ("Event-driven: ", "Functions triggered by HTTP requests, queues, etc."),
        ("Stateless: ", "Functions are ephemeral — no persistent state"),
    ]
    add_bullet_list(slide, Inches(6.7), Inches(2.3), Inches(6.0), Inches(3.0),
                    items, font_size=16, spacing=1.6)

    # Key facts card
    card = add_card(slide, Inches(6.7), Inches(5.5), Inches(5.8), Inches(1.4),
                    fill_color=RGBColor(0xF0, 0xFD, 0xFA))
    add_textbox(slide, Inches(7.0), Inches(5.6), Inches(5.3), Inches(0.5),
                "Two Main Models", font_size=16, color=TEAL, bold=True)
    add_textbox(slide, Inches(7.0), Inches(6.0), Inches(5.3), Inches(0.8),
                "• FaaS (Function-as-a-Service) — deploy code as event-driven functions\n"
                "• BaaS (Backend-as-a-Service) — managed cloud services (DB, storage, auth)",
                font_size=14, color=DARK_TEXT, line_spacing=1.5)


def slide_03_survey(prs):
    """Slide 3: Assigned Paper — Survey on Serverless Computing (deep dive)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_header_bar(slide, "ASSIGNED PAPER: SURVEY ON SERVERLESS COMPUTING", section_num=2)

    # Citation bar
    add_textbox(slide, Inches(0.8), Inches(1.25), Inches(11.5), Inches(0.4),
                "Hassan, Barakat & Sarhan (2021) — Journal of Cloud Computing  •  275 papers reviewed (2016–2020)",
                font_size=15, color=MUTED_TEXT, bold=False)

    # ── LEFT COLUMN: Scope ──
    add_textbox(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(0.5),
                "Scope & Methodology", font_size=22, color=DARK_NAVY, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(2.4), Inches(1.5), Inches(0.04))

    scope_items = [
        ("275 papers ", "systematically reviewed (2016–2020)"),
        ("10 FaaS platforms ", "compared (Lambda, Azure, GCF, OpenWhisk, ...)"),
        ("8 application domains: ", "chatbots, IoT, security, file processing, ..."),
        ("Pricing models: ", "pay-per-invocation, per-duration, per-GB-second"),
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.7), Inches(5.5), Inches(2.5),
                    scope_items, font_size=15, spacing=1.5)

    # History card
    card = add_card(slide, Inches(0.8), Inches(5.0), Inches(5.5), Inches(1.2),
                    fill_color=RGBColor(0xF0, 0xFD, 0xFA))
    add_textbox(slide, Inches(1.1), Inches(5.1), Inches(5.0), Inches(1.0),
                "📅  Timeline:\n"
                "2014 — AWS Lambda launched  •  2016 — Google & Microsoft adopt\n"
                "2020 — Mature ecosystem with 10+ commercial platforms",
                font_size=13, color=TEAL)

    # ── RIGHT COLUMN: Challenges ──
    add_textbox(slide, Inches(7.0), Inches(1.9), Inches(5.5), Inches(0.5),
                "Key Challenges Identified", font_size=22, color=DARK_NAVY, bold=True)
    add_accent_bar(slide, Inches(7.0), Inches(2.4), Inches(1.5), Inches(0.04), color=CORAL)

    challenges = [
        ("Performance ", "(33 papers) — scheduling & call overhead"),
        ("Cold Start Latency ", "(17 papers) — delay initializing idle functions"),
        ("Security ", "(13 papers) — isolation & trust in shared infra"),
        ("Vendor Lock-in ", "— strong dependency on provider ecosystem"),
        ("Limited Duration ", "— short execution time limits"),
        ("Stateless Design ", "— no persistent state across invocations"),
    ]
    add_bullet_list(slide, Inches(7.0), Inches(2.7), Inches(5.5), Inches(3.0),
                    challenges, font_size=15, spacing=1.45, bullet_color=CORAL)

    # Gap callout
    card2 = add_card(slide, Inches(7.0), Inches(5.8), Inches(5.5), Inches(1.0),
                     fill_color=RGBColor(0xFF, 0xF7, 0xED))
    add_textbox(slide, Inches(7.3), Inches(5.9), Inches(5.0), Inches(0.8),
                "⚠️  Critical Gap: No coverage of AI/ML/LLM workloads.\n"
                "      This gap motivates the next two papers.",
                font_size=14, color=RGBColor(0x92, 0x40, 0x0E), bold=True)


def slide_04_roadmap(prs):
    """Slide 4: Two Research Directions — Roadmap with bidirectional figure"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_header_bar(slide, "TWO RESEARCH DIRECTIONS", section_num=3)

    # Subtitle
    add_textbox(slide, Inches(0.8), Inches(1.25), Inches(11.5), Inches(0.4),
                "The survey reveals a gap: no coverage of AI/LLM workloads. Recent research fills this gap in two directions.",
                font_size=16, color=MUTED_TEXT)

    # Bidirectional diagram — center
    bidi_path = os.path.join(FIGURES_DIR, "bidirectional_relationship.png")
    if os.path.exists(bidi_path):
        slide.shapes.add_picture(bidi_path,
                                 Inches(3.5), Inches(1.8), Inches(6.3), Inches(4.2))

    # Left: Direction 1 card
    card1 = add_card(slide, Inches(0.5), Inches(2.0), Inches(3.0), Inches(3.8),
                     fill_color=RGBColor(0xF0, 0xFD, 0xFA))
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0.5), Inches(2.0), Inches(3.0), Inches(0.06), fill_color=TEAL)
    add_textbox(slide, Inches(0.7), Inches(2.2), Inches(2.6), Inches(0.5),
                "Direction 1", font_size=14, color=TEAL, bold=True)
    add_textbox(slide, Inches(0.7), Inches(2.6), Inches(2.6), Inches(0.5),
                "Serverless Serves LLMs", font_size=18, color=DARK_NAVY, bold=True)
    add_textbox(slide, Inches(0.7), Inches(3.2), Inches(2.6), Inches(2.3),
                "How do we deploy LLMs\n"
                "on serverless platforms?\n\n"
                "• Cold start solutions\n"
                "• Multi-tier storage\n"
                "• Live migration\n\n"
                "→ ServerlessLLM (OSDI '24)",
                font_size=13, color=DARK_TEXT, line_spacing=1.4)

    # Right: Direction 2 card
    card2 = add_card(slide, Inches(9.8), Inches(2.0), Inches(3.0), Inches(3.8),
                     fill_color=RGBColor(0xFF, 0xF7, 0xED))
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(9.8), Inches(2.0), Inches(3.0), Inches(0.06), fill_color=CORAL)
    add_textbox(slide, Inches(10.0), Inches(2.2), Inches(2.6), Inches(0.5),
                "Direction 2", font_size=14, color=CORAL, bold=True)
    add_textbox(slide, Inches(10.0), Inches(2.6), Inches(2.6), Inches(0.5),
                "LLMs Improve Serverless", font_size=18, color=DARK_NAVY, bold=True)
    add_textbox(slide, Inches(10.0), Inches(3.2), Inches(2.6), Inches(2.3),
                "How can LLMs make\n"
                "serverless more reliable?\n\n"
                "• Misconfiguration detection\n"
                "• Zero-shot prompting\n"
                "• Chain-of-Thought reasoning\n\n"
                "→ SlsDetector (TOSEM '26)",
                font_size=13, color=DARK_TEXT, line_spacing=1.4)

    # Bottom: presentation structure note
    card3 = add_card(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.9),
                     fill_color=LIGHT_GRAY)
    add_textbox(slide, Inches(1.1), Inches(6.4), Inches(11.2), Inches(0.7),
                "📋  Presentation Structure:   Survey Foundations  →  Direction 1: ServerlessLLM  →  Direction 2: SlsDetector  →  Open Questions",
                font_size=15, color=DARK_NAVY, bold=True)


def slide_05_new_frontier(prs):
    """Slide 5: Why LLMs on Serverless — Benefits & Challenges"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_header_bar(slide, "DIRECTION 1: LLM INFERENCE ON SERVERLESS")

    # Why this matters
    add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.5),
                "Why deploy LLMs on serverless platforms?", font_size=22,
                color=DARK_NAVY, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.9), Inches(2.0), Inches(0.04))

    # Three benefit cards
    benefits = [
        ("Elastic Scaling", "LLM workloads are extremely\nbursty — serverless scales\nautomatically to handle spikes", TEAL),
        ("Cost Efficiency", "Pay only for active inference\ntime — no cost for idle GPUs\nsitting between requests", DARK_BLUE),
        ("Operational Simplicity", "No GPU cluster management\nfor developers — focus on\nthe application logic only", CORAL),
    ]

    for i, (title, desc, color) in enumerate(benefits):
        x = Inches(0.8) + Inches(i * 4.1)
        card = add_card(slide, x, Inches(2.2), Inches(3.8), Inches(2.2),
                        fill_color=LIGHT_GRAY)
        add_shape(slide, MSO_SHAPE.RECTANGLE,
                  x, Inches(2.2), Inches(3.8), Inches(0.06), fill_color=color)
        add_textbox(slide, x + Inches(0.3), Inches(2.4), Inches(3.2), Inches(0.5),
                    title, font_size=20, color=color, bold=True)
        add_textbox(slide, x + Inches(0.3), Inches(2.95), Inches(3.2), Inches(1.2),
                    desc, font_size=15, color=DARK_TEXT, line_spacing=1.4)

    # But: the challenges
    add_textbox(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.5),
                "But: LLMs break fundamental serverless assumptions", font_size=22,
                color=CORAL, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(5.3), Inches(2.5), Inches(0.04), color=CORAL)

    challenge_items = [
        ("Cold start is catastrophic: ", "Loading a 70B-param model takes minutes — "
         "traditional serverless cold start is milliseconds"),
        ("Stateless ≠ LLM inference: ", "KV-cache, conversation history, and multi-turn "
         "reasoning all require state across invocations"),
        ("GPU resource mismatch: ", "Serverless platforms designed for CPU functions — "
         "GPU provisioning is coarse-grained and expensive"),
    ]
    add_bullet_list(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(2.0),
                    challenge_items, font_size=15, spacing=1.5, bullet_color=CORAL)

    # Footnote: mention the other survey
    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.3),
                "See also: Wang et al. (WoSC '24) — first survey of AI model inference on serverless, categorizing ML/DL/LLM workloads [4]",
                font_size=11, color=MUTED_TEXT)


def slide_06_serverlessllm(prs):
    """Slide 6: ServerlessLLM Deep Dive"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_header_bar(slide, "SERVERLESSLLM: SOLVING THE COLD START PROBLEM")

    # Citation
    add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.4),
                "Fu et al. (2024)  —  OSDI '24 (USENIX)  •  One of the premier systems conferences",
                font_size=15, color=MUTED_TEXT)

    # Left: architecture diagram
    arch_path = os.path.join(FIGURES_DIR, "serverlessllm_architecture.png")
    if os.path.exists(arch_path):
        slide.shapes.add_picture(arch_path,
                                 Inches(0.3), Inches(1.7), Inches(5.8), Inches(5.5))

    # Right: Three contributions with detail
    add_textbox(slide, Inches(6.5), Inches(1.7), Inches(6), Inches(0.5),
                "Three Core Innovations", font_size=22, color=DARK_NAVY, bold=True)
    add_accent_bar(slide, Inches(6.5), Inches(2.2), Inches(1.8), Inches(0.04))

    # Contribution 1
    add_textbox(slide, Inches(6.5), Inches(2.5), Inches(6), Inches(0.4),
                "① Multi-Tier Checkpoint Loading", font_size=17, color=TEAL, bold=True)
    add_bullet_list(slide, Inches(6.8), Inches(2.9), Inches(5.7), Inches(1.0),
                    ["Loading-optimized checkpoint format (sequential, chunk-based)",
                     "Exploits full storage hierarchy: GPU → DRAM → SSD → remote",
                     "Result: 3.6–8.2× faster than PyTorch/Safetensors"],
                    font_size=13, spacing=1.3, bullet_color=TEAL)

    # Contribution 2
    add_textbox(slide, Inches(6.5), Inches(4.1), Inches(6), Inches(0.4),
                "② LLM Live Migration (first ever)", font_size=17, color=TEAL, bold=True)
    add_bullet_list(slide, Inches(6.8), Inches(4.5), Inches(5.7), Inches(0.8),
                    ["Migrates only tokens — NOT the massive KV-cache",
                     "Re-computes KV-cache at destination, cutting network traffic"],
                    font_size=13, spacing=1.3, bullet_color=TEAL)

    # Contribution 3
    add_textbox(slide, Inches(6.5), Inches(5.4), Inches(6), Inches(0.4),
                "③ Locality-Aware Scheduling (Phantom)", font_size=17, color=TEAL, bold=True)
    add_bullet_list(slide, Inches(6.8), Inches(5.8), Inches(5.7), Inches(0.8),
                    ["Cost models estimate load time per storage tier",
                     "Picks optimal server to minimize startup latency"],
                    font_size=13, spacing=1.3, bullet_color=TEAL)

    # Bottom result card
    card = add_card(slide, Inches(6.5), Inches(6.5), Inches(6.0), Inches(0.7),
                    fill_color=RGBColor(0xF0, 0xFD, 0xFA))
    add_textbox(slide, Inches(6.8), Inches(6.55), Inches(5.5), Inches(0.6),
                "🚀  10–200× latency improvement over KServe & Ray Serve (Azure Trace workloads)",
                font_size=15, color=TEAL, bold=True)


def slide_07_serverlessllm_eval(prs):
    """Slide 7: ServerlessLLM Evaluation Details"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_header_bar(slide, "SERVERLESSLLM: EVALUATION & IMPACT")

    # Citation
    add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.4),
                "Fu et al. (2024)  —  OSDI '24  •  Evaluated on GPU clusters with real-world workloads",
                font_size=15, color=MUTED_TEXT)

    # Left: Evaluation setup
    add_textbox(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(0.5),
                "Evaluation Setup", font_size=22, color=DARK_NAVY, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(2.4), Inches(1.5), Inches(0.04))

    setup = [
        ("Hardware: ", "8× A5000 GPUs, 1TB DDR4, NVMe SSDs (RAID 0)"),
        ("Cluster: ", "4 servers × 4 A40 GPUs, 10Gbps Ethernet"),
        ("Models: ", "OPT (2.7B–66B), LLaMA-2 (7B–70B), Falcon (7B–40B)"),
        ("Datasets: ", "GSM8K (math reasoning), ShareGPT (chat logs)"),
        ("Workloads: ", "Azure Serverless Trace — bursty, real-world patterns"),
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.7), Inches(5.5), Inches(2.8),
                    setup, font_size=14, spacing=1.6)

    # Right: Key results
    add_textbox(slide, Inches(7.0), Inches(1.9), Inches(5.5), Inches(0.5),
                "Key Results", font_size=22, color=DARK_NAVY, bold=True)
    add_accent_bar(slide, Inches(7.0), Inches(2.4), Inches(1.5), Inches(0.04), color=TEAL)

    # Result metric cards
    result_data = [
        ("3.6–8.2×", "Faster checkpoint\nloading"),
        ("10–200×", "Latency improvement\nover baselines"),
        ("4.4×", "LoRA adaptor\nloading speedup"),
    ]
    for i, (value, label) in enumerate(result_data):
        x = Inches(7.0) + Inches(i * 2.0)
        card = add_card(slide, x, Inches(2.7), Inches(1.8), Inches(1.8),
                        fill_color=LIGHT_GRAY)
        add_textbox(slide, x + Inches(0.1), Inches(2.8), Inches(1.6), Inches(0.7),
                    value, font_size=26, color=TEAL, bold=True,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.1), Inches(3.4), Inches(1.6), Inches(0.7),
                    label, font_size=12, color=DARK_TEXT,
                    alignment=PP_ALIGN.CENTER, line_spacing=1.4)

    # Why it matters
    add_textbox(slide, Inches(7.0), Inches(4.8), Inches(5.5), Inches(0.5),
                "Why This Matters", font_size=20, color=DARK_NAVY, bold=True)
    add_accent_bar(slide, Inches(7.0), Inches(5.25), Inches(1.5), Inches(0.04))

    impact = [
        "Proves serverless LLM inference is practical at production scale",
        "Multi-tier storage design is generalizable to any GPU cluster",
        "Live migration enables load balancing without dropping requests",
        "Directly addresses cold start — the #1 challenge from the survey",
    ]
    add_bullet_list(slide, Inches(7.0), Inches(5.5), Inches(5.5), Inches(2.0),
                    impact, font_size=14, spacing=1.4, bullet_color=TEAL)

    # Brief mention of other work
    card = add_card(slide, Inches(0.8), Inches(5.7), Inches(5.5), Inches(1.5),
                    fill_color=RGBColor(0xF8, 0xFA, 0xFC))
    add_textbox(slide, Inches(1.1), Inches(5.8), Inches(5.0), Inches(0.4),
                "Related Recent Work", font_size=16, color=MUTED_TEXT, bold=True)
    add_textbox(slide, Inches(1.1), Inches(6.2), Inches(5.0), Inches(0.9),
                "• SLINFER (HPCA '26): adds CPU+GPU heterogeneous sharing,\n"
                "  achieving 86–154% capacity improvement [6]\n"
                "• Samanta & Nguyen (WoSC '25): proposes 3-layer architecture\n"
                "  vision for unified serverless LLM systems [5]",
                font_size=12, color=MUTED_TEXT, line_spacing=1.4)


def slide_08_slsdetector(prs):
    """Slide 8: SlsDetector — LLMs Improving Serverless (full-width layout)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_header_bar(slide, "DIRECTION 2: LLMs IMPROVING SERVERLESS")

    # Citation
    add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.4),
                "Wen et al. (2026)  —  ACM TOSEM  •  SlsDetector: LLM-Based Misconfiguration Detection for AWS SAM",
                font_size=15, color=MUTED_TEXT)

    # ── LEFT COLUMN: Problem ──
    add_textbox(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.5),
                "The Misconfiguration Problem", font_size=22, color=DARK_NAVY, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(2.3), Inches(1.8), Inches(0.04), color=CORAL)

    problems = [
        ("800+ AWS resource types ", "with complex dependency chains"),
        ("50,000+ IDs exposed ", "from a single S3 bucket misconfiguration"),
        ("4.9M customers breached ", "from API misconfiguration (DoorDash)"),
        ("Traditional methods fail: ", "data-driven approaches need large labeled datasets"),
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(2.3),
                    problems, font_size=15, spacing=1.6, bullet_color=CORAL)

    # ── LEFT COLUMN: Approach ──
    add_textbox(slide, Inches(0.8), Inches(4.7), Inches(5.5), Inches(0.5),
                "SlsDetector: The LLM Approach", font_size=20, color=DARK_NAVY, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(5.15), Inches(1.8), Inches(0.04), color=TEAL)

    approach = [
        ("Zero-shot prompting ", "— no training data required at all"),
        ("Chain-of-Thought ", "— step-by-step reasoning for complex configs"),
        ("Multi-dim constraints ", "— aligned to serverless config semantics"),
        ("Structured output ", "— deterministic, explainable detection results"),
    ]
    add_bullet_list(slide, Inches(0.8), Inches(5.4), Inches(5.5), Inches(2.0),
                    approach, font_size=15, spacing=1.5)

    # ── RIGHT COLUMN: Results ──
    add_textbox(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(0.5),
                "Results (ChatGPT-4o on 110 config files)", font_size=22,
                color=DARK_NAVY, bold=True)
    add_accent_bar(slide, Inches(7.0), Inches(2.3), Inches(1.8), Inches(0.04), color=TEAL)

    # Result metric cards — larger
    metrics = [
        ("72.88%", "Precision", "+53.82 pp"),
        ("88.18%", "Recall", "+17.40 pp"),
        ("79.75%", "F1-Score", "+49.72 pp"),
    ]
    for i, (value, label, improvement) in enumerate(metrics):
        x = Inches(7.0) + Inches(i * 2.0)
        card = add_card(slide, x, Inches(2.6), Inches(1.8), Inches(2.2),
                        fill_color=LIGHT_GRAY)
        add_textbox(slide, x + Inches(0.1), Inches(2.7), Inches(1.6), Inches(0.7),
                    value, font_size=30, color=TEAL, bold=True,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.1), Inches(3.4), Inches(1.6), Inches(0.4),
                    label, font_size=16, color=DARK_NAVY, bold=True,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.1), Inches(3.8), Inches(1.6), Inches(0.4),
                    improvement + " vs. SOTA", font_size=13, color=TEAL,
                    alignment=PP_ALIGN.CENTER)

    # Generalization section
    add_textbox(slide, Inches(7.0), Inches(5.1), Inches(5.5), Inches(0.5),
                "Generalizes Across LLMs", font_size=20, color=DARK_NAVY, bold=True)
    add_accent_bar(slide, Inches(7.0), Inches(5.55), Inches(1.5), Inches(0.04), color=TEAL)

    llms = [
        "ChatGPT-4o  •  Llama 3.1 (405B) Instruct",
        "Gemini 1.5 Pro  •  DeepSeek V3",
        "Consistent high effectiveness across all tested models",
    ]
    add_bullet_list(slide, Inches(7.0), Inches(5.8), Inches(5.5), Inches(1.5),
                    llms, font_size=14, spacing=1.5, bullet_color=TEAL)


def slide_09_synthesis_and_questions(prs):
    """Slide 9: Synthesis + Open Research Question"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_header_bar(slide, "SYNTHESIS & OPEN RESEARCH QUESTION")

    # ── LEFT: Narrative arc summary ──
    add_textbox(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5),
                "The Story So Far", font_size=22, color=DARK_NAVY, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.9), Inches(1.8), Inches(0.04))

    # Narrative steps with connecting flow
    steps = [
        ("Survey (2021):", " Serverless is mature, but cold start,\n"
         "  statelessness, and security are unsolved.", DARK_BLUE),
        ("Gap:", " No coverage of AI/LLM workloads — the next\n"
         "  frontier for serverless computing.", CORAL),
        ("ServerlessLLM (2024):", " Solves cold start with\n"
         "  multi-tier loading (10–200× improvement).", TEAL),
        ("SlsDetector (2026):", " LLMs can improve serverless\n"
         "  itself — bidirectional relationship.", AMBER),
    ]

    y = 2.2
    for title, desc, color in steps:
        add_shape(slide, MSO_SHAPE.RECTANGLE,
                  Inches(0.8), Inches(y), Inches(0.12), Inches(1.0), fill_color=color)
        txBox = slide.shapes.add_textbox(Inches(1.15), Inches(y), Inches(5.0), Inches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        run1 = tf.paragraphs[0].add_run()
        run1.text = title
        run1.font.size = Pt(15)
        run1.font.color.rgb = color
        run1.font.bold = True
        run1.font.name = "Calibri"
        run2 = tf.paragraphs[0].add_run()
        run2.text = desc
        run2.font.size = Pt(14)
        run2.font.color.rgb = DARK_TEXT
        run2.font.bold = False
        run2.font.name = "Calibri"
        tf.paragraphs[0].line_spacing = Pt(20)
        y += 1.25

    # ── RIGHT: Open Research Question ──
    add_textbox(slide, Inches(7.0), Inches(1.4), Inches(5.5), Inches(0.5),
                "Open Research Question", font_size=22, color=DARK_NAVY, bold=True)
    add_accent_bar(slide, Inches(7.0), Inches(1.9), Inches(1.8), Inches(0.04), color=CORAL)

    # The question
    card = add_card(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(2.0),
                    fill_color=RGBColor(0xFF, 0xF7, 0xED))
    add_textbox(slide, Inches(7.3), Inches(2.3), Inches(5.0), Inches(0.6),
                "Can LLM agents autonomously manage\nserverless infrastructure?",
                font_size=20, color=CORAL, bold=True, line_spacing=1.3)
    add_textbox(slide, Inches(7.3), Inches(3.1), Inches(5.0), Inches(0.9),
                "Extending SlsDetector's approach from static detection\n"
                "to real-time autonomous infrastructure management.",
                font_size=14, color=DARK_TEXT, line_spacing=1.4)

    # Why this question
    add_textbox(slide, Inches(7.0), Inches(4.5), Inches(5.5), Inches(0.4),
                "Why This Question?", font_size=18, color=DARK_NAVY, bold=True)
    add_accent_bar(slide, Inches(7.0), Inches(4.9), Inches(1.2), Inches(0.04), color=TEAL)

    reasons = [
        "SlsDetector proves LLMs understand serverless configs",
        "Auto-scaling, resource allocation, scheduling are all\n  decision problems LLMs could reason about",
        "Would close the loop: LLMs serving on serverless,\n  managed by LLMs — a self-optimizing system",
        "Combines cloud computing + AI — core CS 6620 themes",
    ]
    add_bullet_list(slide, Inches(7.0), Inches(5.1), Inches(5.5), Inches(2.0),
                    reasons, font_size=13, spacing=1.45, bullet_color=TEAL)


def slide_10_conclusion(prs):
    """Slide 10: Conclusion & Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)

    # Accent bars
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), fill_color=TEAL)
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06), fill_color=TEAL)

    # Left accent
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(1.0), Inches(1.2), Inches(0.06), Inches(4.0), fill_color=TEAL)

    # Title
    add_textbox(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(0.8),
                "Key Takeaways", font_size=36, color=WHITE, bold=True)

    takeaways = [
        ("1. ", "Serverless computing is evolving from general FaaS to AI-specific platforms"),
        ("2. ", "Cold start is the #1 barrier — ServerlessLLM achieves 10–200× improvement"),
        ("3. ", "The relationship is bidirectional: serverless serves LLMs, LLMs improve serverless"),
        ("4. ", "Open question: Can LLM agents autonomously manage serverless infrastructure?"),
    ]

    y = 2.3
    for num, text in takeaways:
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(y), Inches(10), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        run1 = tf.paragraphs[0].add_run()
        run1.text = num
        run1.font.size = Pt(18)
        run1.font.color.rgb = LIGHT_TEAL
        run1.font.bold = True
        run1.font.name = "Calibri"
        run2 = tf.paragraphs[0].add_run()
        run2.text = text
        run2.font.size = Pt(17)
        run2.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
        run2.font.bold = False
        run2.font.name = "Calibri"
        y += 0.7

    # Divider
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(1.5), Inches(5.3), Inches(3), Inches(0.03), fill_color=TEAL)

    # Thank you
    add_textbox(slide, Inches(1.5), Inches(5.6), Inches(10), Inches(0.8),
                "Thank You — Questions?", font_size=32, color=LIGHT_TEAL, bold=True)

    add_textbox(slide, Inches(1.5), Inches(6.4), Inches(10), Inches(0.5),
                "CS 6620 — Cloud Computing  •  Spring 2026", font_size=16,
                color=MUTED_TEXT)


def slide_11_references(prs):
    """Slide 11: References (backup)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_header_bar(slide, "REFERENCES")

    refs = [
        ("[1]  ", 'Hassan, Barakat & Sarhan. "Survey on Serverless Computing."\n'
                  "      Journal of Cloud Computing, 10(39), 2021.  [Assigned Paper]", True),
        ("[2]  ", 'Fu et al. "ServerlessLLM: Low-Latency Serverless Inference for LLMs."\n'
                  "      OSDI '24, USENIX, pp. 135-153, 2024.", True),
        ("[3]  ", 'Wen et al. "LLM-Based Misconfiguration Detection for AWS Serverless Computing."\n'
                  "      ACM Trans. Softw. Eng. Methodol., 35(4), Article 110, 2026.", True),
        ("[4]  ", 'Wang, Jiang & Mi. "Advancing Serverless Computing for Scalable AI Model Inference."\n'
                  "      WoSC '24, ACM, pp. 1-6, 2024.", False),
        ("[5]  ", 'Samanta & Nguyen. "Illuminating the Hidden Challenges of Serverless LLM Systems."\n'
                  "      WoSC '25, ACM, pp. 51-57, 2025.", False),
        ("[6]  ", 'Xu et al. "Towards Resource-Efficient Serverless LLM Inference with SLINFER."\n'
                  "      HPCA '26, IEEE, pp. 1-18, 2026.", False),
    ]

    y = 1.5
    for num, text, is_core in refs:
        color = TEAL if is_core else MUTED_TEXT
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.5), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        run1 = tf.paragraphs[0].add_run()
        run1.text = num
        run1.font.size = Pt(14)
        run1.font.color.rgb = color
        run1.font.bold = True
        run1.font.name = "Calibri"
        run2 = tf.paragraphs[0].add_run()
        run2.text = text
        run2.font.size = Pt(14)
        run2.font.color.rgb = DARK_TEXT if is_core else MUTED_TEXT
        run2.font.bold = is_core
        run2.font.name = "Calibri"
        tf.paragraphs[0].line_spacing = Pt(18)
        y += 0.9

    # Legend
    add_textbox(slide, Inches(0.8), Inches(y + 0.3), Inches(11.5), Inches(0.4),
                "Papers [1]–[3] are the core papers presented in depth.  "
                "Papers [4]–[6] are supplementary references mentioned briefly.",
                font_size=12, color=MUTED_TEXT)


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Build all slides (11 total)
    slide_01_title(prs)                    # 1: Title
    slide_02_background(prs)               # 2: Background - Serverless
    slide_03_survey(prs)                   # 3: Assigned Paper (deep)
    slide_04_roadmap(prs)                  # 4: Two Research Directions (bidirectional fig)
    slide_05_new_frontier(prs)             # 5: Direction 1 - LLM on Serverless
    slide_06_serverlessllm(prs)            # 6: ServerlessLLM (deep)
    slide_07_serverlessllm_eval(prs)       # 7: ServerlessLLM Evaluation
    slide_08_slsdetector(prs)              # 8: Direction 2 - SlsDetector (deep)
    slide_09_synthesis_and_questions(prs)   # 9: Synthesis + Open Question
    slide_10_conclusion(prs)               # 10: Conclusion & Thank You
    slide_11_references(prs)               # 11: References (backup)

    prs.save(OUTPUT_FILE)
    print(f"✅ Presentation saved to: {OUTPUT_FILE}")
    print(f"   Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
